"""上传文件到纯文本知识文档的解析器。"""

import csv
import json
from dataclasses import dataclass
from io import BytesIO, StringIO
from pathlib import PurePosixPath
from typing import Any


class UnsupportedDocumentFormat(ValueError):
    pass


class DocumentParseError(ValueError):
    pass


@dataclass
class ParsedDocument:
    title: str
    content: str
    metadata: dict[str, Any]


class DocumentIngestionService:
    """根据文件扩展名调用相应解析器，并统一返回可索引的文本。"""

    max_extracted_characters = 100_000
    supported_extensions = {
        ".md",
        ".markdown",
        ".txt",
        ".csv",
        ".tsv",
        ".json",
        ".pdf",
        ".docx",
        ".pptx",
        ".xlsx",
        ".xlsm",
        ".xls",
    }

    def parse(
        self, filename: str, payload: bytes, content_type: str | None = None
    ) -> ParsedDocument:
        """解析单个文件。

        先限制格式、空文件和提取文本长度，再交给文档仓库。二进制原件不会直接发送给模型。
        """
        extension = PurePosixPath(filename).suffix.lower()
        if extension not in self.supported_extensions:
            supported = ", ".join(sorted(self.supported_extensions - {".xlsm"}))
            raise UnsupportedDocumentFormat(
                f"unsupported file format '{extension or 'unknown'}'; supported formats: {supported}"
            )
        if not payload:
            raise DocumentParseError("uploaded file is empty")
        title = (
            PurePosixPath(filename).stem.replace("_", " ").replace("-", " ").strip() or "未命名文档"
        )
        metadata = {
            "source_filename": filename,
            "source_content_type": content_type or "application/octet-stream",
            "parser": extension.lstrip("."),
        }
        try:
            if extension in {".md", ".markdown", ".txt"}:
                content = payload.decode("utf-8-sig")
            elif extension in {".csv", ".tsv"}:
                content, rows = self._parse_delimited(payload, "\t" if extension == ".tsv" else ",")
                metadata["row_count"] = rows
            elif extension == ".json":
                content = json.dumps(
                    json.loads(payload.decode("utf-8-sig")), ensure_ascii=False, indent=2
                )
            elif extension == ".pdf":
                content, pages = self._parse_pdf(payload)
                metadata["page_count"] = pages
            elif extension == ".docx":
                content, paragraphs, tables = self._parse_docx(payload)
                metadata["paragraph_count"] = paragraphs
                metadata["table_count"] = tables
            elif extension == ".pptx":
                content, slides = self._parse_pptx(payload)
                metadata["slide_count"] = slides
            else:
                content, sheets = self._parse_excel(payload, extension)
                metadata["sheet_count"] = sheets
        except Exception as exc:
            raise DocumentParseError(f"could not parse '{filename}': {exc}") from exc
        if not content.strip():
            raise DocumentParseError(f"'{filename}' contains no extractable text")
        if len(content) > self.max_extracted_characters:
            raise DocumentParseError(
                f"'{filename}' exceeds the {self.max_extracted_characters} character extraction limit"
            )
        metadata["character_count"] = len(content)
        return ParsedDocument(title=title, content=content.strip(), metadata=metadata)

    @staticmethod
    def _parse_delimited(payload: bytes, delimiter: str) -> tuple[str, int]:
        text = payload.decode("utf-8-sig")
        rows = list(csv.reader(StringIO(text), delimiter=delimiter))
        return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows if any(row)), len(
            rows
        )

    @staticmethod
    def _parse_pdf(payload: bytes) -> tuple[str, int]:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(payload))
        pages = [page.extract_text() or "" for page in reader.pages]
        return "\n\n".join(
            f"[第 {index} 页]\n{text}" for index, text in enumerate(pages, 1) if text.strip()
        ), len(pages)

    @staticmethod
    def _parse_docx(payload: bytes) -> tuple[str, int, int]:
        from docx import Document

        document = Document(BytesIO(payload))
        parts = [
            paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
        ]
        for table_index, table in enumerate(document.tables, 1):
            rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
            if rows:
                parts.append(f"[表格 {table_index}]\n" + "\n".join(rows))
        return "\n\n".join(parts), len(document.paragraphs), len(document.tables)

    @staticmethod
    def _parse_pptx(payload: bytes) -> tuple[str, int]:
        from pptx import Presentation

        presentation = Presentation(BytesIO(payload))
        slides: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, 1):
            texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            if texts:
                slides.append(f"[第 {slide_index} 页]\n" + "\n".join(texts))
        return "\n\n".join(slides), len(presentation.slides)

    @staticmethod
    def _parse_excel(payload: bytes, extension: str) -> tuple[str, int]:
        if extension == ".xls":
            import xlrd

            workbook = xlrd.open_workbook(file_contents=payload)
            sheets = [
                f"## {sheet.name}\n"
                + "\n".join(
                    " | ".join(str(value).strip() for value in sheet.row_values(row_index))
                    for row_index in range(sheet.nrows)
                )
                for sheet in workbook.sheets()
            ]
            return "\n\n".join(sheets), workbook.nsheets

        from openpyxl import load_workbook

        workbook = load_workbook(BytesIO(payload), read_only=True, data_only=True)
        sheets: list[str] = []
        for worksheet in workbook.worksheets:
            rows = []
            for row in worksheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value).strip() for value in row]
                if any(values):
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"## {worksheet.title}\n" + "\n".join(rows))
        return "\n\n".join(sheets), len(workbook.worksheets)
